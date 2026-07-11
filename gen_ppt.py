import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, r1, g1, b1, r2, g2, b2):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r1, g1, b1)

def add_shape(slide, left, top, width, height, r, g, b, alpha=None, shape_type=MSO_SHAPE.RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    shape.line.fill.background()
    if alpha:
        shape.fill.fore_color.brightness = 0
    return shape

def add_rounded_rect(slide, left, top, width, height, r, g, b):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text="", font_size=18, color=(255,255,255), bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = RGBColor(*color)
    p.font.bold = bold
    p.alignment = alignment
    return txBox

def add_accent_bar(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*color)
    shape.line.fill.background()
    return shape

# ============ 封面 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, 15, 23, 42, 20, 40, 80)

# 顶部装饰条 - 渐变效果用多个色块模拟
colors_top = [(0,112,243), (0,180,216), (144,224,239)]
for i, c in enumerate(colors_top):
    add_accent_bar(slide, Inches(i*4.444), Inches(0), Inches(4.444), Inches(0.08), c)

# 左侧大装饰块
add_shape(slide, Inches(0), Inches(0), Inches(0.6), Inches(7.5), 0, 112, 243)

# 底部装饰色块
add_shape(slide, Inches(0), Inches(6.8), Inches(13.333), Inches(0.7), 0, 180, 216)

# 主标题
add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.5), 
    "项 目 申 报 书", font_size=54, color=(255,255,255), bold=True, alignment=PP_ALIGN.LEFT)

# 装饰线
add_accent_bar(slide, Inches(1.5), Inches(3.7), Inches(2), Inches(0.06), (0, 180, 216))

# 副标题
add_text_box(slide, Inches(1.5), Inches(4.0), Inches(10), Inches(0.8),
    "项 目 名 称  ·  申 报 单 位  ·  日 期", font_size=24, color=(144, 224, 239), bold=False, alignment=PP_ALIGN.LEFT)

# 右下角装饰圆
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(4.5), Inches(2.5), Inches(2.5))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0, 112, 243)
circle.fill.fore_color.brightness = 0.0
circle.line.fill.background()
# 设置透明度 - 用RGBA方式
from pptx.oxml.ns import qn
solidFill = circle.fill._fill
solidFill.set(qn('a:alpha'), '30000')

# ============ 目录页 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, 245, 247, 250, 240, 242, 248)

# 左侧色块
add_shape(slide, Inches(0), Inches(0), Inches(0.8), Inches(7.5), 0, 112, 243)

# 顶部装饰
add_accent_bar(slide, Inches(0.8), Inches(0), Inches(12.533), Inches(0.06), (0, 112, 243))

# 目录标题
add_text_box(slide, Inches(1.5), Inches(0.8), Inches(5), Inches(0.6),
    "目 录", font_size=40, color=(15, 23, 42), bold=True)

# 目录项
items = ["一、项目背景与意义", "二、项目目标与定位", "三、项目核心内容",
         "四、技术路线与方案", "五、预期成果与效益", "六、实施计划与进度", "七、经费预算与保障"]
for i, item in enumerate(items):
    y = Inches(1.8) + Inches(i * 0.72)
    # 编号圆
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), y, Inches(0.45), Inches(0.45))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0, 112, 243)
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = str(i+1)
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    add_text_box(slide, Inches(2.2), y, Inches(8), Inches(0.45),
        item, font_size=20, color=(30, 30, 40))

# 右侧装饰
add_shape(slide, Inches(10.5), Inches(1.5), Inches(2.5), Inches(5.5), 0, 112, 243, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)

chapters = [
    ("项目背景", "背景/背景/背景"),
    ("项目目标", "目标/目标/目标"),
    ("核心内容", "内容/内容/内容"),
    ("技术方案", "方案/方案/方案"),
    ("预期成果", "成果/成果/成果"),
    ("实施计划", "计划/计划/计划"),
    ("经费预算", "预算/预算/预算")
]

chapter_colors = [
    (0, 112, 243),
    (0, 180, 216),
    (72, 149, 239),
    (0, 119, 182),
    (2, 62, 138),
    (0, 150, 199),
    (33, 158, 188)
]

# 每章2页 = 分隔页 + 内容页
for idx, (ch, desc) in enumerate(chapters):
    c = chapter_colors[idx]
    
    # === 分隔页 ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, c[0], c[1], c[2], c[0], c[1], c[2])
    
    add_text_box(slide, Inches(1.5), Inches(2.5), Inches(2), Inches(0.6),
        f"0{idx+1}", font_size=72, color=(255,255,255, 30), bold=True)
    
    add_accent_bar(slide, Inches(1.5), Inches(3.5), Inches(3), Inches(0.08), (255,255,255))
    
    add_text_box(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(1.2),
        ch, font_size=44, color=(255,255,255), bold=True)
    
    add_text_box(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(0.6),
        f"第{idx+1}章", font_size=18, color=(200,220,255), bold=False)
    
    # 装饰圆
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(1.5), Inches(3.5), Inches(3.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(255,255,255)
    circle.fill.fore_color.brightness = 0.0
    circle.line.fill.background()
    solidFill = circle.fill._fill
    solidFill.set(qn('a:alpha'), '10000')
    
    # === 内容页 ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, 245, 247, 250, 240, 242, 248)
    
    # 顶部条
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.12), c[0], c[1], c[2])
    
    # 左侧装饰条
    add_shape(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), c[0], c[1], c[2])
    
    # 章节标题
    add_text_box(slide, Inches(1.2), Inches(0.5), Inches(10), Inches(0.5),
        f"0{idx+1}  {ch}", font_size=28, color=(c[0], c[1], c[2]), bold=True)
    
    # 装饰线
    add_accent_bar(slide, Inches(1.2), Inches(1.1), Inches(2), Inches(0.04), (c[0], c[1], c[2]))
    
    # 内容卡片 - 三个卡片
    for card in range(3):
        x = Inches(1.2) + Inches(card * 3.8)
        y = Inches(1.6)
        w = Inches(3.4)
        h = Inches(3.8)
        
        card_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        card_shape.fill.solid()
        card_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card_shape.line.color.rgb = RGBColor(220, 225, 235)
        card_shape.line.width = Pt(1)
        
        # 卡片顶部色条
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.06))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(c[0], c[1], c[2])
        bar.line.fill.background()
    
    # 底部装饰
    add_shape(slide, Inches(0), Inches(6.8), Inches(13.333), Inches(0.7), 240, 242, 248)

# ============ 尾页 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, 15, 23, 42, 20, 40, 80)

add_shape(slide, Inches(0), Inches(0), Inches(0.6), Inches(7.5), 0, 112, 243)
add_shape(slide, Inches(0), Inches(6.8), Inches(13.333), Inches(0.7), 0, 180, 216)

add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1.0),
    "感 谢 关 注", font_size=48, color=(255,255,255), bold=True, alignment=PP_ALIGN.CENTER)

add_accent_bar(slide, Inches(5.5), Inches(3.7), Inches(2.333), Inches(0.06), (0, 180, 216))

add_text_box(slide, Inches(1.5), Inches(4.0), Inches(10), Inches(0.8),
    "联系人：____________    电话：____________    邮箱：____________",
    font_size=20, color=(144, 224, 239), alignment=PP_ALIGN.CENTER)

output_path = os.path.expanduser("~/Downloads/无敌了/青云/项目申报书_高级版.pptx")
prs.save(output_path)
print(f"✅ PPT已生成: {output_path}")
